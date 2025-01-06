import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from transformers import T5ForConditionalGeneration, T5Tokenizer
import genanki
from pptx import Presentation
import random
import logging
import openai
import time
import os

logging.basicConfig(level=logging.INFO)

openai.api_key = 'YOUR_API_KEY'  

def select_ppt_file():
    file_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Select a PowerPoint file"
    )
    if not file_path:
        messagebox.showerror("Error", "No file selected")
        return None
    return file_path

def select_save_location(default_filename):
    file_path = filedialog.asksaveasfilename(
        defaultextension=".apkg",
        filetypes=[("Anki Deck files", "*.apkg")],
        title="Save Anki Deck As",
        initialfile=default_filename
    )
    if not file_path:
        messagebox.showerror("Error", "No save location selected")
        return None
    return file_path

def get_presentation_title(ppt):
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip()
                title = title.split('\n')[0]
                title = title.replace(':', '-')
                return title[:50]
    return "Anki_Deck"

def generate_anki_deck(ppt_path, progress_var, time_var):
    ppt = Presentation(ppt_path)
    deck_title = get_presentation_title(ppt)
    deck_id = random.getrandbits(31)
    model_id = random.getrandbits(31)

    my_deck = genanki.Deck(deck_id=deck_id, name=deck_title)

    basic_model = genanki.Model(
        model_id=model_id,
        name='Basic Model',
        fields=[{'name': 'Front'}, {'name': 'Back'}],
        templates=[{
            'name': 'Card 1',
            'qfmt': '{{Front}}',
            'afmt': '{{Back}}'
        }],
    )

    def generate_questions_answers(text):
        prompt = (f"Create meaningful and relevant Anki flashcards from the following text. "
                  f"Ensure that questions are specific and directly related to the content, and "
                  f"provide clear, concise answers:\n{text}")
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that generates Anki flashcards."},
                    {"role": "user", "content": prompt}
                ]
            )
            message = response.choices[0].message['content']
            return parse_questions_and_answers(message)
        except Exception as e:
            logging.error(f"Error generating questions and answers with GPT: {e}")
            return []

    def parse_questions_and_answers(gpt_output):
        q_and_a_pairs = []
        questions = gpt_output.split("Question:")
        for q in questions[1:]:
            parts = q.split("Answer:")
            if len(parts) == 2:
                question = parts[0].strip()
                answer = parts[1].strip()
                q_and_a_pairs.append((question, answer))
        return q_and_a_pairs

    def is_meaningful_question(question):
        return len(question) > 5 and '?' in question

    def is_good_answer(answer):
        return len(answer) > 5

    def process_slide_text(slide):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + " "
        return slide_text.strip()

    total_slides = len(ppt.slides)
    start_time = time.time()
    
    for i, slide in enumerate(ppt.slides):
        slide_text = process_slide_text(slide)
        if slide_text:
            qa_pairs = generate_questions_answers(slide_text)
            for question, answer in qa_pairs:
                try:
                    if is_meaningful_question(question) and is_good_answer(answer):
                        note = genanki.Note(
                            model=basic_model,
                            fields=[question, answer],
                            guid=str(random.getrandbits(63))
                        )
                        my_deck.add_note(note)
                except Exception as e:
                    logging.error(f"Error adding note to deck: {e}")
        
        progress_var.set((i + 1) / total_slides * 100)
        elapsed_time = time.time() - start_time
        estimated_time = (elapsed_time / (i + 1)) * (total_slides - (i + 1))
        time_var.set(f"Estimated Time Remaining: {int(estimated_time)} seconds")
        root.update_idletasks()

    output_path = select_save_location(f'{deck_title}.apkg')
    if output_path:
        try:
            genanki.Package(my_deck).write_to_file(output_path)
            messagebox.showinfo("Success", f"Anki deck saved to {output_path}")
        except Exception as e:
            logging.error(f"Error saving Anki deck: {e}")
            messagebox.showerror("Error", f"Error saving Anki deck: {e}")

def on_generate_anki_deck():
    ppt_path = select_ppt_file()
    if ppt_path:
        progress_var.set(0)
        time_var.set("Estimated Time Remaining: Calculating...")
        generate_anki_deck(ppt_path, progress_var, time_var)

root = tk.Tk()
root.title("Anki Deck Generator")
root.geometry("600x250") 

frame = tk.Frame(root, padx=20, pady=20)
frame.pack(fill='both', expand=True)

progress_var = tk.DoubleVar()
progress_bar = ttk.Progressbar(frame, variable=progress_var, maximum=100, length=550)
progress_bar.pack(pady=10)

time_var = tk.StringVar()
time_label = tk.Label(frame, textvariable=time_var, font=('Arial', 12))
time_label.pack(pady=5)

button_frame = tk.Frame(frame)
button_frame.pack(pady=10)

tk.Button(button_frame, text="Drop PowerPoint File", command=on_generate_anki_deck).pack(side='left', padx=5)
tk.Button(button_frame, text="Exit", command=root.quit).pack(side='right', padx=5)

root.mainloop()
